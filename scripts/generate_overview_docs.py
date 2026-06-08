#!/usr/bin/env python3
"""Generate Project_Overview.pdf and Project_Overview.pptx (neutral technical explanation)."""

from __future__ import annotations

import os
import subprocess
import sys

DOCS_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs"
)


def build_pptx(path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x1A, 0x73, 0xE8)
    DARK = RGBColor(0x20, 0x21, 0x24)
    GRAY = RGBColor(0x5F, 0x63, 0x68)

    def add_title_slide(title: str, subtitle: str = "") -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = DARK
        if subtitle:
            box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.5))
            tf2 = box2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = GRAY

    def add_slide(title: str, bullets: list[str], note: str = "") -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # title bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLUE
        bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.7))
        tp = tb.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(24)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(5.5))
        tf = body.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
            p.font.color.rgb = DARK
            p.space_after = Pt(10)
            p.level = 0
        if note:
            nb = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.7))
            np = nb.text_frame.paragraphs[0]
            np.text = note
            np.font.size = Pt(14)
            np.font.italic = True
            np.font.color.rgb = GRAY

    add_title_slide(
        "ML Fault Detection & Self-Healing\nin Telecom Networks",
        "Technical overview — simulation, machine learning, and MAPE-K evaluation",
    )

    add_slide(
        "Purpose",
        [
            "End-to-end pipeline for proactive RAN fault detection in dense HetNet networks",
            "Combines NS-3 simulation, ML classifiers (RF, LSTM, SVM), and MAPE-K autonomic control",
            "Research question: Does ML proactive detection reduce MTTR and improve availability vs. reactive threshold monitoring?",
            "Results are observed from simulation — no artificial blending toward target tables by default",
        ],
    )

    add_slide(
        "Real-World Operational Context",
        [
            "Simulation imitates Nigerian MNO operations (e.g. MTN, Airtel) — not a live network integration",
            "Power fault → site outage (grid / generator / rectifier failure)",
            "Congestion → peak-hour or event traffic surge",
            "Hardware fault → RRU/BBU failure or transport cut",
            "Reactive baseline: severe OSS-style thresholds, 240–300 min manual remediation",
            "MAPE-K path: earlier ML detection, 35–55 min modelled automated recovery",
        ],
    )

    add_slide(
        "System Architecture",
        [
            "Stage 1: NS-3 HetNet simulation → per-cell KPI CSV files",
            "Stage 2: Merge trials → kpi_master_dataset.csv",
            "Stage 3: Sliding windows → train RF, LSTM, SVM",
            "Stage 4: MAPE-K loop on test data → MTTR and availability",
            "MAPE-K: Monitor → Analyse (ML) → Plan → Execute",
        ],
    )

    add_slide(
        "Network Topology (Table 3.1)",
        [
            "7 macro + 21 small cells = 28 cells total",
            "Hexagonal layout, inter-site distance 500 m",
            "500 UEs (design target); 280 UEs stable default on NS-3",
            "300 s simulation; 1 s KPI interval",
            "50 trials × 4 fault types = 200 simulation runs",
            "Full dataset: ~1.68 million raw CSV rows",
        ],
    )

    add_slide(
        "Three Simulation Backends",
        [
            "thesis-fault-sim — fast KPI generator for large-scale ML datasets",
            "thesis-fault-sim-lte — LENA LTE/EPC with real PHY traces + FlowMonitor",
            "thesis-fault-sim-nr — 5G-LENA NR/EPC (n78, 3.5 GHz) for 5G validation",
            "All backends use identical CSV schema and fault labelling",
            "Commands: run_all_trials.py | --lte | --nr",
        ],
    )

    add_slide(
        "Fault Classes",
        [
            "0 — Normal: no active fault window",
            "1 — Power fault: gNB transmit power collapse",
            "2 — Congestion: UDP traffic surge on affected cell UEs",
            "3 — Hardware failure: gNB PHY deactivation",
            "Stochastic per trial: random cell, onset time, duration",
            "Ground truth stored in fault_label column for supervised ML",
        ],
    )

    add_slide(
        "KPI Features (8 per cell per second)",
        [
            "RSRP, SINR, PRB utilisation",
            "Downlink and uplink throughput",
            "Packet loss rate, handover success rate, latency",
            "Collected every 1 s for each of 28 cells",
        ],
    )

    add_slide(
        "Machine Learning Pipeline",
        [
            "10-second sliding windows, stride 1, per cell",
            "48 tabular features (6 statistics × 8 KPIs) + sequences for LSTM",
            "Train / val / test: 70 / 15 / 15 % stratified split",
            "SMOTE on training data only; PCA 95% variance on tabular features",
            "Models: Random Forest (tabular), LSTM (temporal), SVM (baseline)",
            "Output: reports/ml_metrics.json (observed accuracy, F1, AUC)",
        ],
    )

    add_slide(
        "MAPE-K Evaluation",
        [
            "Monitor: KPI window pre-filter (PRB, RSRP, loss, throughput)",
            "Analyse: ML inference with confidence threshold 0.70",
            "Plan / Execute: fault-specific remediation (Table 4.6, minutes)",
            "Compared against reactive baseline (severe thresholds only)",
            "Outputs: MTTR and network availability per model",
            "NCC 99% availability benchmark in Chapter 4 figures",
        ],
        note="Reference targets: Reactive ~312 min MTTR; LSTM ~102 min; availability ~99%",
    )

    add_slide(
        "Dataset Tracks",
        [
            "Track A — KPI generator: full scale, fast, simplified RAN physics",
            "Track B — LTE: real LENA traces, LTE-A HetNet surrogate for 5G",
            "Track C — Both: ML/MAPE-K on KPI data + LTE cross-validation",
            "Document which backend produced each result set",
        ],
    )

    add_slide(
        "Key Outputs & Files",
        [
            "output/kpi_master_dataset.csv — merged labelled dataset",
            "reports/ml_metrics.json — classifier performance",
            "reports/mapek_summary.json — MTTR and availability",
            "reports/fig3_*.png — Chapter 3 methodology figures",
            "reports/fig4_*.png — Chapter 4 results figures",
            "thesis_constants.py — approved parameters and reference targets",
        ],
    )

    add_slide(
        "Limitations",
        [
            "Simulation study — not live commercial RAN deployment",
            "MAPE-K Execute models recovery time; no real OSS API integration",
            "280 UE default due to NS-3 stability (500 is design target)",
            "Some Chapter 3 figures are schematic methodology diagrams",
            "NR trials require significantly longer runtime than KPI generator",
            "Metric blend weight = 0 (observed-only reporting)",
        ],
    )

    add_slide(
        "Technology Stack",
        [
            "NS-3 3.38 — network simulator",
            "LENA (LTE) + 5G-LENA v2.4 (NR) — RAN stacks",
            "TensorFlow/Keras 2.15 + scikit-learn 1.5 — ML",
            "imbalanced-learn — SMOTE",
            "Pipeline: setup.sh → run_all_trials.py → preprocess_and_train.py → mapek_loop.py",
        ],
    )

    prs.save(path)
    print(f"  Wrote {path}")


def build_html(path: str) -> None:
    html = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>ML Fault Detection & Self-Healing — Technical Overview</title>
<style>
  @page { margin: 18mm 16mm; size: A4; }
  * { box-sizing: border-box; }
  body {
    font-family: "Segoe UI", Helvetica, Arial, sans-serif;
    font-size: 10.5pt;
    line-height: 1.45;
    color: #202124;
    max-width: 210mm;
    margin: 0 auto;
    padding: 12mm 14mm;
  }
  h1 { font-size: 22pt; color: #1a73e8; border-bottom: 3px solid #1a73e8; padding-bottom: 8px; margin-top: 0; }
  h2 { font-size: 14pt; color: #1a73e8; margin-top: 22px; page-break-after: avoid; }
  h3 { font-size: 11pt; color: #3c4043; margin-top: 14px; }
  p, li { margin: 6px 0; }
  ul { padding-left: 20px; }
  table { width: 100%; border-collapse: collapse; margin: 10px 0 16px; font-size: 9.5pt; }
  th { background: #1a73e8; color: #fff; text-align: left; padding: 7px 9px; }
  td { border: 1px solid #dadce0; padding: 6px 9px; vertical-align: top; }
  tr:nth-child(even) td { background: #f8f9fa; }
  .subtitle { font-size: 12pt; color: #5f6368; margin-bottom: 20px; }
  .box { background: #e8f0fe; border-left: 4px solid #1a73e8; padding: 10px 14px; margin: 12px 0; }
  pre { background: #f1f3f4; padding: 10px; font-size: 8.5pt; overflow-x: auto; border-radius: 4px; }
  .footer { margin-top: 24px; font-size: 9pt; color: #80868b; border-top: 1px solid #dadce0; padding-top: 10px; }
  .page-break { page-break-before: always; }
</style>
</head>
<body>

<h1>ML Fault Detection &amp; Self-Healing in Telecom Networks</h1>
<p class="subtitle">Technical Overview — Simulation, Machine Learning, and MAPE-K Evaluation</p>

<h2>1. Purpose</h2>
<p>This project implements an end-to-end pipeline for <strong>proactive RAN fault detection</strong> and <strong>MAPE-K-guided self-healing</strong> in dense HetNet mobile networks.</p>
<div class="box"><strong>Research question:</strong> Can machine-learning-based proactive fault detection reduce Mean Time to Recovery (MTTR) and improve network availability compared to conventional threshold-based reactive monitoring?</div>
<p>The pipeline combines NS-3 network simulation, ML classifiers (Random Forest, LSTM, SVM), and a MAPE-K autonomic evaluation loop. Results are produced from <strong>observed simulation outputs</strong> — metric blending toward draft tables is disabled by default.</p>

<h2>2. Real-World Operational Context</h2>
<p>This is a <strong>simulation study</strong> that imitates operational scenarios common in Nigerian mobile networks (e.g. MTN, Airtel). It is not a live operator network integration.</p>
<table>
<tr><th>Operational scenario</th><th>Field cause</th><th>Simulation proxy</th><th>KPI effect</th></tr>
<tr><td>Cell outage</td><td>Grid failure, generator fault</td><td>Power fault</td><td>RSRP ↓, throughput → 0</td></tr>
<tr><td>Slow / congested data</td><td>Peak-hour surge</td><td>Congestion</td><td>PRB &gt; 90%, latency ↑</td></tr>
<tr><td>Site / HW failure</td><td>RRU/BBU, fibre cut</td><td>HW fault</td><td>Partial collapse, high loss</td></tr>
<tr><td>Normal</td><td>Typical urban traffic</td><td>None campaign</td><td>Stable KPIs</td></tr>
</table>
<p><strong>Reactive model:</strong> severe thresholds only → 240–300 min manual remediation.<br>
<strong>MAPE-K model:</strong> earlier ML detection → 35–55 min automated recovery (Table 4.6).<br>
<strong>Regulatory reference:</strong> NCC 99% availability benchmark in Chapter 4 figures.</p>

<h2>3. System Architecture</h2>
<pre>NS-3 HetNet Simulation  →  KPI CSV Dataset  →  ML Training  →  MAPE-K Evaluation
                                                              →  MTTR + Availability</pre>
<p><strong>MAPE-K loop:</strong> Monitor (KPI window) → Analyse (ML, confidence ≥ 0.70) → Plan (remediation policy) → Execute (recovery time).</p>

<div class="page-break"></div>
<h2>4. Network Topology</h2>
<table>
<tr><th>Parameter</th><th>Value</th></tr>
<tr><td>Macro cells</td><td>7 (hexagonal, ISD 500 m)</td></tr>
<tr><td>Small cells</td><td>21 (3 per macro)</td></tr>
<tr><td>Total cells</td><td><strong>28</strong></td></tr>
<tr><td>UEs</td><td>500 target / 280 stable default</td></tr>
<tr><td>Simulation time</td><td>300 s (120 s in batch runs)</td></tr>
<tr><td>Trials</td><td>50 × 4 faults = 200 runs</td></tr>
<tr><td>Raw dataset size</td><td>~1,680,000 rows (full merge)</td></tr>
<tr><td>ML windows (after subsampling)</td><td>~51,000</td></tr>
</table>

<h2>5. Three Simulation Backends</h2>
<table>
<tr><th>Backend</th><th>Script</th><th>Speed</th><th>Use case</th></tr>
<tr><td>KPI generator</td><td>thesis-fault-sim</td><td>Fastest</td><td>Large-scale ML dataset</td></tr>
<tr><td>LTE HetNet</td><td>thesis-fault-sim-lte</td><td>Slow</td><td>Real PHY traces (LENA/EPC)</td></tr>
<tr><td>5G NR HetNet</td><td>thesis-fault-sim-nr</td><td>Slowest</td><td>True NR (5G-LENA n78)</td></tr>
</table>
<pre>python3 run_all_trials.py
python3 run_all_trials.py --lte --sim-time 120 --num-ues 280
python3 run_all_trials.py --nr  --sim-time 120 --num-ues 280 --workers 1</pre>

<h2>6. Fault Classes &amp; KPI Features</h2>
<table>
<tr><th>Label</th><th>Class</th><th>Mechanism</th></tr>
<tr><td>0</td><td>Normal</td><td>No fault window</td></tr>
<tr><td>1</td><td>Power fault</td><td>gNB TX power collapse</td></tr>
<tr><td>2</td><td>Congestion</td><td>Traffic surge on cell UEs</td></tr>
<tr><td>3</td><td>Hardware failure</td><td>PHY deactivation</td></tr>
</table>
<p><strong>8 KPIs per cell per second:</strong> RSRP, SINR, PRB utilisation, DL/UL throughput, packet loss, handover success rate, latency.</p>

<h2>7. Machine Learning Pipeline</h2>
<ul>
<li>10-second sliding windows, stride 1, per cell</li>
<li>48 tabular features + LSTM sequences</li>
<li>70 / 15 / 15 % stratified split; SMOTE on train only</li>
<li>Models: Random Forest, LSTM, SVM</li>
<li>Output: <code>reports/ml_metrics.json</code></li>
</ul>

<h2>8. MAPE-K Evaluation</h2>
<table>
<tr><th>Approach</th><th>MTTR (min)</th><th>Availability (%)</th></tr>
<tr><td>Reactive baseline</td><td>~312</td><td>~94.2</td></tr>
<tr><td>LSTM + MAPE-K</td><td>~102</td><td>~99.0</td></tr>
<tr><td>RF + MAPE-K</td><td>~119</td><td>~98.1</td></tr>
<tr><td>SVM + MAPE-K</td><td>~187</td><td>~96.7</td></tr>
</table>
<p><em>Report observed values from JSON artefacts. Values above are reference targets for comparison.</em></p>

<div class="page-break"></div>
<h2>9. Dataset Tracks</h2>
<table>
<tr><th>Track</th><th>Backend</th><th>Strength</th><th>Limitation</th></tr>
<tr><td>A</td><td>KPI generator only</td><td>Full scale, fast</td><td>Simplified RAN physics</td></tr>
<tr><td>B</td><td>LTE only</td><td>Real LENA traces</td><td>LTE surrogate; 280 UEs default</td></tr>
<tr><td>C</td><td>KPI + LTE</td><td>Scale + RAN validation</td><td>State which data source per result</td></tr>
</table>

<h2>10. Key Outputs</h2>
<ul>
<li><code>output/kpi_master_dataset.csv</code> — merged labelled dataset</li>
<li><code>reports/ml_metrics.json</code> — classifier metrics</li>
<li><code>reports/mapek_summary.json</code> — MTTR and availability</li>
<li><code>reports/fig3_*.png</code>, <code>reports/fig4_*.png</code> — thesis figures</li>
</ul>

<h2>11. Limitations</h2>
<ul>
<li>Simulation — not live commercial RAN</li>
<li>MAPE-K Execute models recovery; no real OSS APIs</li>
<li>280 UE stability default (500 is design target)</li>
<li>NR requires long runtime</li>
<li>Some Ch. 3 figures are schematic</li>
</ul>

<h2>12. Technology Stack</h2>
<p>NS-3 3.38 · LENA (LTE) · 5G-LENA v2.4 (NR) · TensorFlow/Keras · scikit-learn · SMOTE</p>

<div class="footer">ML Fault Detection &amp; Self-Healing in Telecom Networks — Technical Overview Document</div>
</body>
</html>"""
    with open(path, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"  Wrote {path}")


def html_to_pdf(html_path: str, pdf_path: str) -> None:
    chrome = "/usr/bin/google-chrome"
    if not os.path.isfile(chrome):
        for c in ("chromium", "chromium-browser", "google-chrome-stable"):
            if os.path.isfile(f"/usr/bin/{c}"):
                chrome = f"/usr/bin/{c}"
                break
    cmd = [
        chrome,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        f"--print-to-pdf={pdf_path}",
        html_path,
    ]
    try:
        subprocess.run(cmd, check=True, capture_output=True, timeout=60)
        print(f"  Wrote {pdf_path}")
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as e:
        print(f"  PDF via Chrome failed: {e}")
        print("  Open docs/Project_Overview.html in a browser → Print → Save as PDF")


def main() -> int:
    os.makedirs(DOCS_DIR, exist_ok=True)
    pptx_path = os.path.join(DOCS_DIR, "Project_Overview.pptx")
    html_path = os.path.join(DOCS_DIR, "Project_Overview.html")
    pdf_path = os.path.join(DOCS_DIR, "Project_Overview.pdf")

    print("Generating overview documents...")
    build_html(html_path)
    build_pptx(pptx_path)
    html_to_pdf(html_path, pdf_path)
    print("Done.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
